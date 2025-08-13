from flask import Flask, request, jsonify, send_file
import os
import jinja2
import uuid
import logging
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('ppt_service_pptx.log')
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

app.config['OUTPUT_FOLDER'] = 'static/pptx'

# 确保目录存在
def ensure_directories():
    logger.info("确保目录存在...")
    try:
        os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
        logger.info(f"目录已创建或存在: {app.config['OUTPUT_FOLDER']}")
    except Exception as e:
        logger.error(f"创建目录失败 {app.config['OUTPUT_FOLDER']}: {str(e)}")

# 模板环境
try:
    env = jinja2.Environment(loader=jinja2.FileSystemLoader('templates'))
    logger.info("Jinja2模板环境初始化成功")
except Exception as e:
    logger.error(f"模板环境初始化失败: {str(e)}")
    env = None

def create_editable_ppt(title, content, author, date):
    """创建可编辑的PPTX文件"""
    try:
        # 创建新的演示文稿
        prs = Presentation()
        
        # 设置幻灯片尺寸为16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 第一页：标题页
        title_slide_layout = prs.slide_layouts[0]  # 标题页布局
        slide = prs.slides.add_slide(title_slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        # 设置副标题
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"作者: {author}\n日期: {date}"
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
        
        # 目录页
        toc_slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        toc_slide = prs.slides.add_slide(toc_slide_layout)
        
        toc_slide.shapes.title.text = "目录"
        toc_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        toc_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        
        # 添加目录内容
        toc_body = toc_slide.placeholders[1]
        toc_text = ""
        for i, section in enumerate(content.get('sections', []), 1):
            toc_text += f"{i}. {section.get('title', '未命名章节')}\n"
        
        toc_body.text = toc_text
        toc_body.text_frame.paragraphs[0].font.size = Pt(24)
        
        # 为每个章节创建幻灯片
        for i, section in enumerate(content.get('sections', []), 1):
            # 章节标题页
            section_title_slide = prs.slides.add_slide(toc_slide_layout)
            section_title_slide.shapes.title.text = f"第{i}章: {section.get('title', '未命名章节')}"
            section_title_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
            section_title_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            
            # 章节内容页
            content_slide = prs.slides.add_slide(toc_slide_layout)
            content_slide.shapes.title.text = section.get('title', '未命名章节')
            content_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
            content_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            
            # 添加章节描述
            content_body = content_slide.placeholders[1]
            content_body.text = section.get('description', '暂无描述')
            content_body.text_frame.paragraphs[0].font.size = Pt(20)
        
        # 结束页
        end_slide = prs.slides.add_slide(toc_slide_layout)
        end_slide.shapes.title.text = "谢谢观看"
        end_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        end_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        end_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        end_body = end_slide.placeholders[1]
        end_body.text = f"作者: {author}\n日期: {date}"
        end_body.text_frame.paragraphs[0].font.size = Pt(24)
        end_body.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
        
        logger.info(f"成功创建PPT，共 {len(prs.slides)} 张幻灯片")
        return prs
        
    except Exception as e:
        logger.error(f"创建PPT失败: {str(e)}")
        return None

def generate_ppt_with_pptx(title, content, author, date, output_path):
    """使用python-pptx生成可编辑的PPTX"""
    try:
        logger.info("使用python-pptx创建可编辑PPT...")
        
        # 创建PPT
        prs = create_editable_ppt(title, content, author, date)
        if not prs:
            return False
        
        # 保存文件
        prs.save(output_path)
        
        # 验证文件
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"PPT文件已保存: {output_path}, 大小: {file_size} 字节")
            
            if file_size < 1000:
                logger.error(f"文件过小 ({file_size} 字节)")
                return False
            
            # 验证文件格式
            try:
                with open(output_path, 'rb') as f:
                    header = f.read(4)
                    if header != b'PK\x03\x04':
                        logger.error(f"文件格式错误: {repr(header)}")
                        return False
                    logger.info("文件格式验证通过（ZIP格式）")
            except Exception as e:
                logger.error(f"文件验证失败: {str(e)}")
                return False
            
            return True
        else:
            logger.error("文件保存失败")
            return False
            
    except Exception as e:
        logger.error(f"生成PPT失败: {str(e)}")
        return False

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    """PPT生成接口"""
    logger.info("收到PPT生成请求")
    
    try:
        # 确保目录存在
        ensure_directories()
        
        if env is None:
            logger.error("模板环境未初始化")
            return jsonify({"error": "模板环境未初始化"}), 500
        
        data = request.json
        logger.info(f"接收到的数据: {data}")
        
        # 验证必要参数
        if not data or 'title' not in data:
            logger.warning("缺少必要参数: title")
            return jsonify({"error": "缺少必要参数: title"}), 400
        
        # 验证content结构
        content = data.get('content', {})
        if not content or 'sections' not in content:
            logger.warning("缺少content.sections参数")
            return jsonify({"error": "缺少content.sections参数"}), 400
        
        sections = content.get('sections', [])
        if not sections:
            logger.warning("sections数组为空")
            return jsonify({"error": "sections数组不能为空"}), 400
        
        logger.info(f"处理 {len(sections)} 个章节")
        
        # 生成唯一文件名
        pptx_filename = f"lecture_{uuid.uuid4().hex}.pptx"
        pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], pptx_filename)
        logger.info(f"PPT文件路径: {pptx_path}")
        
        # 生成PPT
        if generate_ppt_with_pptx(
            title=data['title'],
            content=content,
            author=data.get('author', '智课云师'),
            date=data.get('date', ''),
            output_path=pptx_path
        ):
            logger.info("PPT生成成功")
            return jsonify({
                "status": "success",
                "download_url": f"/api/download-ppt/{pptx_filename}",
                "filename": pptx_filename,
                "message": "使用python-pptx生成，完全可编辑"
            })
        else:
            logger.error("PPT生成失败")
            return jsonify({"error": "PPT生成失败"}), 500
            
    except Exception as e:
        logger.error(f"PPT生成过程中发生未预期的错误: {str(e)}")
        import traceback
        logger.error(f"错误堆栈: {traceback.format_exc()}")
        return jsonify({"error": f"服务器内部错误: {str(e)}"}), 500

@app.route('/api/download-ppt/<filename>')
def download_ppt(filename):
    """PPT下载接口"""
    logger.info(f"下载请求: {filename}")
    pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if os.path.exists(pptx_path):
        return send_file(pptx_path, as_attachment=True)
    else:
        logger.warning(f"文件不存在: {pptx_path}")
        return jsonify({"error": "文件不存在"}), 404

@app.route('/api/health')
def health_check():
    """健康检查接口"""
    return jsonify({
        "status": "healthy", 
        "service": "PPT Generator (python-pptx)",
        "features": ["完全可编辑", "标准PPTX格式", "兼容所有Office软件"]
    })

@app.route('/')
def index():
    """首页"""
    return """
    <h1>智能课程PPT生成服务 (python-pptx版本)</h1>
    <p>✅ 生成完全可编辑的PPTX文件</p>
    <p>✅ 兼容PowerPoint、WPS等所有Office软件</p>
    <p>✅ 支持添加、删除、修改幻灯片</p>
    <p>使用 POST /api/generate-ppt 接口生成PPT</p>
    <p>示例请求:</p>
    <pre>
    {
        "title": "测试课程",
        "content": {
            "sections": [
                {
                    "title": "第一章节",
                    "description": "这是第一个测试章节的内容"
                }
            ]
        }
    }
    </pre>
    <p>查看<a href="/api/health">服务状态</a></p>
    """

if __name__ == '__main__':
    logger.info("PPT生成服务正在启动 (python-pptx版本)...")
    logger.info(f"服务地址: http://localhost:5000")
    logger.info("✅ 生成完全可编辑的PPTX文件")
    ensure_directories()
    app.run(host='0.0.0.0', port=5000, debug=True) 