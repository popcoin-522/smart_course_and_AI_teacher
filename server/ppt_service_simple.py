from flask import Flask, request, jsonify, send_file
import os
import uuid
import logging
from flask_cors import CORS

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('ppt_service_simple.log')
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

app.config['OUTPUT_FOLDER'] = 'static/pptx'

# 确保目录存在
def ensure_directories():
    logger.info("确保目录存在...")
    try:
        os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
        logger.info(f"目录已创建或存在: {app.config['OUTPUT_FOLDER']}")
    except Exception as e:
        logger.error(f"创建目录失败 {app.config['OUTPUT_FOLDER']}: {str(e)}")

def create_simple_ppt(title, content, author, date):
    """创建简单的可编辑PPTX文件"""
    try:
        # 导入放在函数内部，避免启动时的导入错误
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        logger.info("开始创建PPT...")
        
        # 创建新的演示文稿
        prs = Presentation()
        logger.info("演示文稿对象创建成功")
        
        # 第一页：标题页
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        logger.info(f"标题页创建成功: {title}")
        
        # 设置副标题
        if hasattr(slide.placeholders, '__getitem__'):
            try:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = f"作者: {author}\n日期: {date}"
                logger.info("副标题设置成功")
            except Exception as e:
                logger.warning(f"副标题设置失败: {e}")
        
        # 目录页
        toc_slide_layout = prs.slide_layouts[1]
        toc_slide = prs.slides.add_slide(toc_slide_layout)
        
        toc_slide.shapes.title.text = "目录"
        logger.info("目录页创建成功")
        
        # 添加目录内容
        try:
            toc_body = toc_slide.placeholders[1]
            toc_text = ""
            for i, section in enumerate(content.get('sections', []), 1):
                toc_text += f"{i}. {section.get('title', '未命名章节')}\n"
            
            toc_body.text = toc_text
            logger.info("目录内容设置成功")
        except Exception as e:
            logger.warning(f"目录内容设置失败: {e}")
        
        # 为每个章节创建幻灯片
        sections = content.get('sections', [])
        logger.info(f"开始创建 {len(sections)} 个章节的幻灯片")
        
        for i, section in enumerate(sections, 1):
            try:
                # 章节标题页
                section_title_slide = prs.slides.add_slide(toc_slide_layout)
                section_title_slide.shapes.title.text = f"第{i}章: {section.get('title', '未命名章节')}"
                logger.info(f"章节 {i} 标题页创建成功")
                
                # 章节内容页
                content_slide = prs.slides.add_slide(toc_slide_layout)
                content_slide.shapes.title.text = section.get('title', '未命名章节')
                
                # 添加章节描述
                try:
                    content_body = content_slide.placeholders[1]
                    content_body.text = section.get('description', '暂无描述')
                    logger.info(f"章节 {i} 内容页创建成功")
                except Exception as e:
                    logger.warning(f"章节 {i} 内容设置失败: {e}")
                    
            except Exception as e:
                logger.error(f"创建章节 {i} 幻灯片失败: {e}")
                continue
        
        # 结束页
        try:
            end_slide = prs.slides.add_slide(toc_slide_layout)
            end_slide.shapes.title.text = "谢谢观看"
            logger.info("结束页创建成功")
        except Exception as e:
            logger.warning(f"结束页创建失败: {e}")
        
        logger.info(f"成功创建PPT，共 {len(prs.slides)} 张幻灯片")
        return prs
        
    except Exception as e:
        logger.error(f"创建PPT失败: {str(e)}")
        import traceback
        logger.error(f"错误堆栈: {traceback.format_exc()}")
        return None

def generate_ppt_simple(title, content, author, date, output_path):
    """使用简化方法生成可编辑的PPTX"""
    try:
        logger.info("使用简化方法创建可编辑PPT...")
        
        # 创建PPT
        prs = create_simple_ppt(title, content, author, date)
        if not prs:
            logger.error("PPT创建失败")
            return False
        
        # 保存文件
        logger.info(f"正在保存文件到: {output_path}")
        prs.save(output_path)
        logger.info("文件保存完成")
        
        # 验证文件
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"PPT文件已保存: {output_path}, 大小: {file_size} 字节")
            
            if file_size < 1000:
                logger.error(f"文件过小 ({file_size} 字节)")
                return False
            
            # 验证文件格式
            try:
                with open(output_path, 'rb') as f:
                    header = f.read(4)
                    if header != b'PK\x03\x04':
                        logger.error(f"文件格式错误: {repr(header)}")
                        return False
                    logger.info("文件格式验证通过（ZIP格式）")
            except Exception as e:
                logger.error(f"文件验证失败: {str(e)}")
                return False
            
            logger.info("PPT生成成功！")
            return True
        else:
            logger.error("文件保存失败")
            return False
            
    except Exception as e:
        logger.error(f"生成PPT失败: {str(e)}")
        import traceback
        logger.error(f"错误堆栈: {traceback.format_exc()}")
        return False

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    """PPT生成接口"""
    logger.info("收到PPT生成请求")
    
    try:
        # 确保目录存在
        ensure_directories()
        
        data = request.json
        logger.info(f"接收到的数据: {data}")
        
        # 验证必要参数
        if not data or 'title' not in data:
            logger.warning("缺少必要参数: title")
            return jsonify({"error": "缺少必要参数: title"}), 400
        
        # 验证content结构
        content = data.get('content', {})
        if not content or 'sections' not in content:
            logger.warning("缺少content.sections参数")
            return jsonify({"error": "缺少content.sections参数"}), 400
        
        sections = content.get('sections', [])
        if not sections:
            logger.warning("sections数组为空")
            return jsonify({"error": "sections数组不能为空"}), 400
        
        logger.info(f"处理 {len(sections)} 个章节")
        
        # 生成唯一文件名
        pptx_filename = f"lecture_{uuid.uuid4().hex}.pptx"
        pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], pptx_filename)
        logger.info(f"PPT文件路径: {pptx_path}")
        
        # 生成PPT
        if generate_ppt_simple(
            title=data['title'],
            content=content,
            author=data.get('author', '智课云师'),
            date=data.get('date', ''),
            output_path=pptx_path
        ):
            logger.info("PPT生成成功")
            return jsonify({
                "status": "success",
                "download_url": f"/api/download-ppt/{pptx_filename}",
                "filename": pptx_filename,
                "message": "使用简化方法生成，完全可编辑"
            })
        else:
            logger.error("PPT生成失败")
            return jsonify({"error": "PPT生成失败"}), 500
            
    except Exception as e:
        logger.error(f"PPT生成过程中发生未预期的错误: {str(e)}")
        import traceback
        logger.error(f"错误堆栈: {traceback.format_exc()}")
        return jsonify({"error": f"服务器内部错误: {str(e)}"}), 500

@app.route('/api/download-ppt/<filename>')
def download_ppt(filename):
    """PPT下载接口"""
    logger.info(f"下载请求: {filename}")
    pptx_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    if os.path.exists(pptx_path):
        return send_file(pptx_path, as_attachment=True)
    else:
        logger.warning(f"文件不存在: {pptx_path}")
        return jsonify({"error": "文件不存在"}), 404

@app.route('/api/health')
def health_check():
    """健康检查接口"""
    return jsonify({
        "status": "healthy", 
        "service": "PPT Generator (Simple Version)",
        "features": ["完全可编辑", "标准PPTX格式", "兼容所有Office软件", "简化设计"]
    })

@app.route('/')
def index():
    """首页"""
    return """
    <h1>智能课程PPT生成服务 (简化版本)</h1>
    <p>✅ 生成完全可编辑的PPTX文件</p>
    <p>✅ 兼容PowerPoint、WPS等所有Office软件</p>
    <p>✅ 支持添加、删除、修改幻灯片</p>
    <p>✅ 简化设计，提高稳定性</p>
    <p>使用 POST /api/generate-ppt 接口生成PPT</p>
    <p>示例请求:</p>
    <pre>
    {
        "title": "测试课程",
        "content": {
            "sections": [
                {
                    "title": "第一章节",
                    "description": "这是第一个测试章节的内容"
                }
            ]
        }
    }
    </pre>
    <p>查看<a href="/api/health">服务状态</a></p>
    """

if __name__ == '__main__':
    logger.info("PPT生成服务正在启动 (简化版本)...")
    logger.info(f"服务地址: http://localhost:5000")
    logger.info("✅ 生成完全可编辑的PPTX文件")
    logger.info("✅ 使用简化设计，提高稳定性")
    ensure_directories()
    app.run(host='0.0.0.0', port=5000, debug=True) 