import os
import uuid
import logging
import json
from datetime import datetime
from flask import Flask, request, jsonify, send_from_directory, render_template_string
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import jinja2

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler('ppt_service.log', encoding='utf-8')
    ]
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# 配置
app.config['UPLOAD_FOLDER'] = 'temp'
app.config['PPTX_FOLDER'] = 'static/pptx'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB

# 初始化Jinja2模板环境
try:
    env = jinja2.Environment(loader=jinja2.FileSystemLoader('templates'))
    logger.info("Jinja2模板环境初始化成功")
except Exception as e:
    logger.error(f"模板环境初始化失败: {str(e)}")
    env = None

def ensure_directories():
    """确保必要的目录存在"""
    directories = [
        app.config['UPLOAD_FOLDER'],
        app.config['PPTX_FOLDER']
    ]
    
    for directory in directories:
        if not os.path.exists(directory):
            os.makedirs(directory)
            logger.info(f"目录已创建: {directory}")
        else:
            logger.info(f"目录已存在: {directory}")

def create_title_slide(prs, title, author, date):
    """创建标题幻灯片"""
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 设置副标题
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = f"作者: {author}\n日期: {date}"
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)

def create_toc_slide(prs, sections):
    """创建目录幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = "目录"
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 设置内容
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for i, section in enumerate(sections, 1):
        p = text_frame.add_paragraph()
        p.text = f"{i}. {section['title']}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(68, 68, 68)
        p.space_after = Pt(12)

def create_content_slide(prs, section_title, section_description):
    """创建内容幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = section_title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 设置内容
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = section_description
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(68, 68, 68)
    p.space_after = Pt(12)

def create_summary_slide(prs, author):
    """创建总结幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = "总结"
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
    
    # 设置内容
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    p1 = text_frame.add_paragraph()
    p1.text = "感谢您的关注！"
    p1.font.size = Pt(28)
    p1.font.color.rgb = RGBColor(68, 68, 68)
    p1.space_after = Pt(20)
    
    p2 = text_frame.add_paragraph()
    p2.text = f"联系方式: {author}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(89, 89, 89)

def generate_ppt_from_content(content_data, output_path):
    """使用python-pptx生成PPTX文件"""
    try:
        # 创建新的演示文稿
        prs = Presentation()
        
        # 设置幻灯片尺寸为16:9
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 创建标题幻灯片
        create_title_slide(prs, content_data['title'], content_data['author'], content_data['date'])
        
        # 创建目录幻灯片
        if content_data['content']['sections']:
            create_toc_slide(prs, content_data['content']['sections'])
        
        # 创建内容幻灯片
        for section in content_data['content']['sections']:
            create_content_slide(prs, section['title'], section['description'])
        
        # 创建总结幻灯片
        create_summary_slide(prs, content_data['author'])
        
        # 保存PPTX文件
        prs.save(output_path)
        logger.info(f"PPTX文件已生成: {output_path}")
        
        # 验证输出文件
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"输出文件大小: {file_size} 字节")
            if file_size < 1000:  # 如果文件太小，可能生成失败
                logger.error(f"输出文件过小 ({file_size} 字节)，可能生成失败")
                return False
            
            logger.info("PPTX文件生成成功")
            return True
        else:
            logger.error("输出文件不存在")
            return False
            
    except Exception as e:
        logger.error(f"PPTX生成失败: {str(e)}")
        return False

@app.route('/')
def index():
    """主页"""
    return """
    <!DOCTYPE html>
    <html>
    <head>
        <title>PPT生成服务</title>
        <meta charset="utf-8">
        <style>
            body { font-family: Arial, sans-serif; margin: 40px; }
            .container { max-width: 800px; margin: 0 auto; }
            .header { background: #f8f9fa; padding: 20px; border-radius: 8px; margin-bottom: 20px; }
            .endpoint { background: #e9ecef; padding: 15px; border-radius: 5px; margin: 10px 0; }
            code { background: #f1f3f4; padding: 2px 6px; border-radius: 3px; }
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <h1>🚀 PPT生成服务</h1>
                <p>使用python-pptx生成可编辑的PPTX文件</p>
            </div>
            
            <h2>📋 API接口</h2>
            
            <div class="endpoint">
                <h3>健康检查</h3>
                <p><code>GET /api/health</code></p>
                <p>检查服务状态</p>
            </div>
            
            <div class="endpoint">
                <h3>生成PPT</h3>
                <p><code>POST /api/generate-ppt</code></p>
                <p>根据提供的内容生成PPTX文件</p>
                
                <h4>请求体格式:</h4>
                <pre><code>{
    "title": "演示标题",
    "content": {
        "sections": [
            {
                "title": "章节标题",
                "description": "章节描述内容"
            }
        ]
    },
    "author": "作者姓名",
    "date": "2025/8/13",
    "template": "lecture.md.jinja"
}</code></pre>
            </div>
            
            <h2>🔧 技术特点</h2>
            <ul>
                <li>✅ 使用python-pptx生成原生PPTX文件</li>
                <li>✅ 支持自定义样式和布局</li>
                <li>✅ 生成的文件完全可编辑</li>
                <li>✅ 支持中文内容</li>
                <li>✅ 自动生成目录和总结页</li>
            </ul>
        </div>
    </body>
    </html>
    """

@app.route('/api/health')
def health_check():
    """健康检查接口"""
    return jsonify({
        "status": "healthy",
        "service": "PPT生成服务",
        "technology": "python-pptx",
        "timestamp": datetime.now().isoformat()
    })

@app.route('/api/generate-ppt', methods=['POST'])
def generate_ppt():
    """PPT生成接口"""
    logger.info("收到PPT生成请求")
    
    try:
        # 确保目录存在
        ensure_directories()
        
        if env is None:
            logger.error("模板环境未初始化")
            return jsonify({"error": "模板环境未初始化"}), 500
        
        data = request.json
        logger.info(f"接收到的数据: {data}")
        
        # 验证必要参数
        if not data or 'title' not in data:
            logger.warning("缺少必要参数: title")
            return jsonify({"error": "缺少必要参数: title"}), 400
        
        # 选择模板
        template_name = data.get('template', 'lecture.md.jinja')
        logger.info(f"使用模板: {template_name}")
        
        # 验证content结构
        content = data.get('content', {})
        if not content or 'sections' not in content:
            logger.warning("缺少content.sections参数")
            return jsonify({"error": "缺少content.sections参数"}), 400
        
        sections = content.get('sections', [])
        if not sections:
            logger.warning("sections数组为空")
            return jsonify({"error": "sections数组不能为空"}), 400
        
        logger.info(f"处理 {len(sections)} 个章节")
        
        # 准备模板数据
        template_data = {
            'title': data['title'],
            'author': data.get('author', '未知作者'),
            'date': data.get('date', datetime.now().strftime('%Y/%m/%d')),
            'content': {'sections': sections}
        }
        
        # 渲染模板
        try:
            template = env.get_template(template_name)
            rendered_content = template.render(**template_data)
            logger.info("模板渲染完成")
        except Exception as e:
            logger.error(f"模板渲染失败: {str(e)}")
            return jsonify({"error": f"模板渲染失败: {str(e)}"}), 500
        
        # 生成PPTX文件
        pptx_filename = f"lecture_{uuid.uuid4().hex}.pptx"
        output_path = os.path.join(app.config['PPTX_FOLDER'], pptx_filename)
        
        logger.info(f"PPT文件路径: {output_path}")
        
        # 使用python-pptx生成PPTX
        if generate_ppt_from_content(template_data, output_path):
            logger.info("PPT生成成功")
            
            # 返回成功响应
            return jsonify({
                "success": True,
                "message": "PPT生成成功",
                "filename": pptx_filename,
                "download_url": f"/download/{pptx_filename}",
                "file_size": os.path.getsize(output_path)
            })
        else:
            logger.error("PPT生成失败")
            return jsonify({"error": "PPT生成失败"}), 500
            
    except Exception as e:
        logger.error(f"PPT生成过程中发生错误: {str(e)}")
        return jsonify({"error": f"服务器内部错误: {str(e)}"}), 500

@app.route('/download/<filename>')
def download_pptx(filename):
    """下载PPTX文件"""
    try:
        return send_from_directory(app.config['PPTX_FOLDER'], filename, as_attachment=True)
    except Exception as e:
        logger.error(f"文件下载失败: {str(e)}")
        return jsonify({"error": "文件下载失败"}), 404

if __name__ == '__main__':
    logger.info("PPT生成服务正在启动...")
    logger.info("服务地址: http://localhost:5000")
    
    # 确保目录存在
    ensure_directories()
    
    app.run(host='0.0.0.0', port=5000, debug=True)